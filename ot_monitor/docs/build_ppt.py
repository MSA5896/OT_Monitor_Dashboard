"""
build_ppt.py — Generate OT Infection Monitoring System presentation.
Run from any directory:  python docs/build_ppt.py
Output: docs/OT_Monitor_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Colours (match dashboard theme) ──────────────────────────────────────────
BG       = RGBColor(0x0A, 0x0E, 0x1A)   # deep navy background
SURFACE  = RGBColor(0x12, 0x18, 0x2A)   # card surface
SURFACE2 = RGBColor(0x1A, 0x22, 0x35)   # lighter surface
CYAN     = RGBColor(0x00, 0xE5, 0xFF)   # primary accent
GREEN    = RGBColor(0x00, 0xE6, 0x76)   # good / normal
AMBER    = RGBColor(0xFF, 0xAB, 0x00)   # warning
RED      = RGBColor(0xFF, 0x52, 0x52)   # alarm
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT2    = RGBColor(0x8A, 0x9B, 0xAA)   # muted text
BORDER   = RGBColor(0x1E, 0x2D, 0x40)

# ── Slide dimensions (widescreen 16:9) ───────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

# ── Screenshot paths ─────────────────────────────────────────────────────────
SHOTS = {
    "monitor":  r"C:\Temp\verify_monitor.png",
    "settings": r"C:\Temp\verify_settings.png",
    "alarms":   r"C:\Temp\verify_alarms.png",
    "history":  r"C:\Temp\verify_history.png",
    "logout":   r"C:\Temp\logout_test.png",
}

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank layout


# ── Helper functions ──────────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(blank)
    # Dark background
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return s


def rect(slide, x, y, w, h, fill, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb


def hline(slide, x, y, w, color=BORDER, thickness=1):
    ln = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Pt(thickness))
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def chip(slide, text, x, y, fill, text_color=BG, size=10):
    """Small coloured badge/chip."""
    w = len(text) * 0.085 + 0.2
    r = rect(slide, x, y, w, 0.28, fill)
    r.fill.solid()
    r.fill.fore_color.rgb = fill
    txt(slide, text, x + 0.05, y + 0.02, w, 0.28,
        size=size, bold=True, color=text_color, align=PP_ALIGN.CENTER)


def slide_header(slide, title, subtitle=""):
    """Standard slide title area with cyan accent bar."""
    rect(slide, 0, 0, 13.33, 0.06, CYAN)
    txt(slide, title,    0.4, 0.15, 9, 0.55, size=28, bold=True,  color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.7, 10, 0.35, size=13, bold=False, color=TEXT2)
    hline(slide, 0.4, 1.0, 12.5, BORDER)


def img(slide, path, x, y, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))


def kpi_card(slide, x, y, label, value, unit, status_color, w=1.9, h=1.1):
    """Mini KPI card matching dashboard style."""
    r = rect(slide, x, y, w, h, SURFACE)
    txt(slide, label, x+0.1, y+0.08, w-0.2, 0.22, size=8,  color=TEXT2)
    txt(slide, value, x+0.1, y+0.28, w-0.2, 0.42, size=22, bold=True, color=status_color)
    txt(slide, unit,  x+0.1, y+0.66, w-0.2, 0.25, size=8,  color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 0.08, CYAN)

# Left accent stripe
rect(s, 0, 0.08, 0.05, 7.42, CYAN)

# Main title
txt(s, "OT INFECTION",     0.6, 1.6,  10, 1.0,  size=52, bold=True,  color=WHITE)
txt(s, "MONITORING SYSTEM",0.6, 2.55, 11, 0.9,  size=46, bold=True,  color=CYAN)
hline(s, 0.6, 3.55, 5.5, CYAN, 2)

txt(s, "Real-time Environmental Monitoring for Operating Theatres",
    0.6, 3.75, 9, 0.45, size=15, color=TEXT2)

txt(s, "Raspberry Pi 4  ·  SCD30  ·  BME280  ·  PMS5003",
    0.6, 4.25, 9, 0.4,  size=13, color=AMBER)

# Feature chips
chip(s, " CO₂ ",         0.6,  5.1, CYAN,  BG,  11)
chip(s, " Temperature ", 1.35, 5.1, GREEN, BG,  11)
chip(s, " Humidity ",    2.65, 5.1, AMBER, BG,  11)
chip(s, " PM2.5/PM10 ",  3.55, 5.1, RED,   WHITE, 11)
chip(s, " Pressure ",    4.65, 5.1, SURFACE2, WHITE, 11)

txt(s, "Compliant: NABH · ASHRAE 170 · WHO IEQ · ISO 14644-1",
    0.6, 5.6, 9, 0.4, size=11, color=TEXT2)

# Right side — dashboard preview
if os.path.exists(SHOTS["logout"]):
    img(s, SHOTS["logout"], 7.2, 1.0, 5.9, 3.8)
    # Overlay border
    border = s.shapes.add_shape(1, Inches(7.2), Inches(1.0), Inches(5.9), Inches(3.8))
    border.fill.background()
    border.line.color.rgb = CYAN
    border.line.width = Pt(1.5)

txt(s, "Intelligent Healthcare  ·  Operating Theatre Safety",
    0.6, 6.9, 12, 0.4, size=10, color=BORDER, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Why Monitor the OT Environment?",
             "Surgical site infections (SSI) are directly linked to uncontrolled environmental parameters")

problems = [
    ("🦠", "Surgical Site Infections (SSI)",
     "3–15% of OT procedures result in SSI. Contaminated air is a leading cause.\nWHO estimates SSIs add 10 extra days to hospital stay."),
    ("🌡️", "Temperature & Humidity Drift",
     "ASHRAE 170 mandates 20–24°C and 40–60% RH. Violations degrade sterility\nand affect staff performance and patient thermoregulation."),
    ("💨", "Particulate Matter (PM2.5 / PM10)",
     "Surgical smoke, skin shedding, and HVAC failure spike PM levels.\nISO 14644-1 Class 7 requires <352,000 particles ≥0.5µm per m³."),
    ("🌿", "CO₂ Accumulation",
     "NIOSH action level: 1000 ppm. Elevated CO₂ indicates insufficient\nventilation — raising infection risk and impairing surgeon alertness."),
]

cols = [(0.35, 2.95), (6.85, 2.95), (0.35, 5.35), (6.85, 5.35)]
for i, (icon, title, body) in enumerate(problems):
    x, y = cols[i]
    r = rect(s, x, y, 6.1, 1.95, SURFACE)
    txt(s, icon + "  " + title, x+0.2, y+0.12, 5.7, 0.38,
        size=14, bold=True, color=CYAN)
    hline(s, x+0.2, y+0.5, 5.7, BORDER)
    txt(s, body, x+0.2, y+0.58, 5.6, 1.2, size=11, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Overview
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "System Architecture Overview",
             "Three-layer design: Physical sensors → FastAPI backend → Web dashboard")

# Layer boxes
layers = [
    (0.35,  1.1, 3.8, 5.8, SURFACE,  CYAN,  "HARDWARE LAYER",
     ["SCD30 — CO₂ + Temp + Humidity",
      "BME280 — Barometric Pressure",
      "PMS5003 — PM1 / PM2.5 / PM10",
      "Raspberry Pi 4 (host)",
      "I²C @ 10kHz + UART0 @ 9600"]),
    (4.65,  1.1, 3.8, 5.8, SURFACE,  AMBER, "BACKEND LAYER",
     ["FastAPI + Uvicorn (port 8001)",
      "HardwareSource (asyncio + threads)",
      "AlarmEngine (hysteresis rules)",
      "SQLite storage (90-day retention)",
      "REST API + WebSocket",
      "HMAC-SHA256 session auth"]),
    (8.95,  1.1, 3.8, 5.8, SURFACE,  GREEN, "DASHBOARD LAYER",
     ["Single-page HTML/CSS/JS app",
      "Served from backend (same-origin)",
      "Monitor · Alarms · History · Settings",
      "Role-based UI (admin / viewer)",
      "Live data via REST polling",
      "Runs in any modern browser"]),
]

for x, y, w, h, bg, accent, title, items in layers:
    rect(s, x, y, w, h, bg)
    rect(s, x, y, w, 0.05, accent)
    txt(s, title, x+0.15, y+0.12, w-0.3, 0.35,
        size=11, bold=True, color=accent)
    hline(s, x+0.15, y+0.48, w-0.3, BORDER)
    for j, item in enumerate(items):
        txt(s, "▸  " + item, x+0.15, y+0.55+j*0.72, w-0.3, 0.5,
            size=10.5, color=WHITE)

# Arrows between layers
for ax in [4.35, 8.65]:
    txt(s, "→", ax, 3.5, 0.45, 0.55, size=28, bold=True,
        color=TEXT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Hardware Components
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Hardware Components",
             "Three sensors covering all key OT environmental parameters")

sensors = [
    ("SCD30", "Sensirion", "CO₂ · Temperature · Humidity",
     CYAN, "I²C @ 0x61",
     [("CO₂ Range",    "400–10,000 ppm"),
      ("Temp Range",   "-40 to +70 °C"),
      ("Humidity",     "0–95 %RH"),
      ("Interface",    "I²C (10 kHz)"),
      ("Address",      "0x61 (fixed)"),
      ("Power",        "3.3V / 19 mA"),]),
    ("BME280", "Bosch Sensortec", "Barometric Pressure",
     AMBER, "I²C @ 0x76",
     [("Pressure Range", "300–1100 hPa"),
      ("Accuracy",       "±1 hPa"),
      ("Temp Range",     "-40 to +85 °C"),
      ("Interface",      "I²C (shared bus)"),
      ("Address",        "0x76 (SDO→GND)"),
      ("Power",          "3.3V / 3.6 µA"),]),
    ("PMS5003", "Plantower", "PM1.0 · PM2.5 · PM10",
     RED, "UART0 @ 9600",
     [("PM Range",    "0–500 µg/m³"),
      ("Resolution",  "1 µg/m³"),
      ("Technology",  "Laser scattering"),
      ("Interface",   "UART0 @ 9600 (GPIO14/15)"),
      ("Port",        "/dev/serial0 → ttyAMA0"),
      ("Power",       "5V / 100 mA"),]),
]

for i, (name, mfr, desc, color, iface, specs) in enumerate(sensors):
    x = 0.35 + i * 4.3
    # Card
    rect(s, x, 1.1, 4.0, 5.85, SURFACE)
    rect(s, x, 1.1, 4.0, 0.05, color)
    # Name
    txt(s, name, x+0.2, 1.18, 3.6, 0.45, size=22, bold=True, color=color)
    txt(s, mfr,  x+0.2, 1.58, 3.6, 0.28, size=10, color=TEXT2)
    txt(s, desc, x+0.2, 1.84, 3.6, 0.28, size=11, bold=True, color=WHITE)
    chip(s, iface, x+0.2, 2.18, SURFACE2, color, 10)
    hline(s, x+0.2, 2.55, 3.6, BORDER)
    # Specs
    for j, (k, v) in enumerate(specs):
        y = 2.68 + j * 0.67
        txt(s, k, x+0.2, y, 1.9, 0.45, size=10, color=TEXT2)
        txt(s, v, x+2.1, y, 1.8, 0.45, size=10, bold=True, color=WHITE)

# RPi at the bottom
rect(s, 0.35, 6.65, 12.6, 0.65, SURFACE2)
txt(s, "Raspberry Pi 4 Model B", 0.55, 6.7,  4, 0.5, size=13, bold=True, color=WHITE)
txt(s, "Host platform · 64-bit ARM Cortex-A72 · 2–8 GB RAM · Ubuntu/RPi OS Bookworm",
    4.5, 6.75, 8.2, 0.45, size=10.5, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Wiring Diagram
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "GPIO Wiring Diagram",
             "All sensors connect to the 40-pin GPIO header — no soldering required")

# RPi board outline
rect(s, 0.35, 1.15, 4.5, 6.1, SURFACE)
txt(s, "Raspberry Pi 4", 0.55, 1.22, 4, 0.38, size=13, bold=True, color=WHITE)
hline(s, 0.55, 1.6, 4.1, BORDER)

pins = [
    ("Pin 1",  "3.3V",           CYAN),
    ("Pin 2",  "5V",             RED),
    ("Pin 3",  "GPIO2 / SDA1",   AMBER),
    ("Pin 5",  "GPIO3 / SCL1",   AMBER),
    ("Pin 6",  "GND",            TEXT2),
    ("Pin 9",  "GND",            TEXT2),
    ("Pin 11", "GPIO17",         GREEN),
    ("Pin 13", "GPIO27",         GREEN),
    ("Pin 8",  "GPIO14 / UART0 TX", RGBColor(0xCE,0x93,0xD8)),
    ("Pin 10", "GPIO15 / UART0 RX", RGBColor(0xCE,0x93,0xD8)),
]

for i, (pin, name, col) in enumerate(pins):
    y = 1.75 + i * 0.5
    chip(s, pin, 0.5, y, col, BG, 9)
    txt(s, name, 1.2, y, 3.4, 0.38, size=10, color=WHITE)

# Connection lines area
rect(s, 5.0, 1.15, 7.9, 6.1, SURFACE)
txt(s, "Sensor Connections", 5.2, 1.22, 5, 0.38, size=13, bold=True, color=WHITE)
hline(s, 5.2, 1.6, 7.5, BORDER)

conns = [
    (CYAN,   "SCD30",   "VCC  → Pin 1  (3.3V)"),
    (CYAN,   "SCD30",   "GND  → Pin 6  (GND)"),
    (CYAN,   "SCD30",   "SDA  → Pin 3  (GPIO2 / SDA1)"),
    (CYAN,   "SCD30",   "SCL  → Pin 5  (GPIO3 / SCL1)"),
    (CYAN,   "SCD30",   "SEL  → Pin 6  (GND)  ← I²C mode select"),
    (AMBER,  "BME280",  "VCC  → Pin 1  (3.3V)"),
    (AMBER,  "BME280",  "GND  → Pin 6  (GND)"),
    (AMBER,  "BME280",  "SDA  → Pin 3  (shared with SCD30)"),
    (AMBER,  "BME280",  "SCL  → Pin 5  (shared with SCD30)"),
    (AMBER,  "BME280",  "SDO  → Pin 6  (GND)  ← addr = 0x76"),
    (RED,    "PMS5003", "VCC  → Pin 2  (5V)"),
    (RED,    "PMS5003", "GND  → Pin 9  (GND)"),
    (RED,    "PMS5003", "SET  → Pin 11 (GPIO17 HIGH)"),
    (RED,    "PMS5003", "TX   → Pin 10 (GPIO15 / UART0 RX)"),
    (RED,    "PMS5003", "RX   → Pin 8  (GPIO14 / UART0 TX)"),
]

for i, (col, sensor, conn) in enumerate(conns):
    if i in (0, 5, 10):  # sensor group header
        hline(s, 5.2, 1.7 + i*0.38, 7.5, BORDER)
    chip(s, sensor, 5.2, 1.72 + i*0.38, col, BG, 8)
    txt(s, conn, 6.35, 1.72 + i*0.38, 6.4, 0.34, size=9.5, color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Dashboard: Monitor Panel
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Dashboard — Monitor Panel",
             "Real-time KPI cards, trend charts, and particulate matter bars updated every 2 seconds")

if os.path.exists(SHOTS["monitor"]):
    img(s, SHOTS["monitor"], 0.35, 1.1, 9.1, 6.15)
    border = s.shapes.add_shape(1, Inches(0.35), Inches(1.1), Inches(9.1), Inches(6.15))
    border.fill.background(); border.line.color.rgb = BORDER; border.line.width = Pt(1)

# Callouts
callouts = [
    (9.6,  1.2, GREEN,  "Live KPI Cards",
     "Temperature, Humidity, PM1, PM2.5,\nPM10, CO₂, Battery, Power Source\nupdated every 2 seconds"),
    (9.6,  2.8, AMBER,  "Battery Warning Banner",
     "Auto-shown when battery < 30%.\nCovers the full width for immediate\nstaff attention"),
    (9.6,  4.4, CYAN,   "Trend Chart",
     "10-minute rolling temperature\nand humidity waveforms"),
    (9.6,  5.8, RED,    "PM Bar Chart",
     "PM1 / PM2.5 / PM10 relative\nbars with real µg/m³ values"),
]

for x, y, col, title, body in callouts:
    rect(s, x, y, 3.5, 1.35, SURFACE)
    rect(s, x, y, 0.04, 1.35, col)
    txt(s, title, x+0.15, y+0.06, 3.2, 0.3, size=11, bold=True, color=col)
    txt(s, body,  x+0.15, y+0.38, 3.2, 0.9, size=9.5, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Dashboard: Settings Panel
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Dashboard — Settings Panel",
             "Admin-only threshold configuration and user management; viewer gets read-only access")

if os.path.exists(SHOTS["settings"]):
    img(s, SHOTS["settings"], 0.35, 1.1, 9.1, 6.15)
    border = s.shapes.add_shape(1, Inches(0.35), Inches(1.1), Inches(9.1), Inches(6.15))
    border.fill.background(); border.line.color.rgb = BORDER; border.line.width = Pt(1)

callouts = [
    (9.6,  1.2, AMBER,  "Sensor Thresholds Table",
     "Per-parameter Ideal / Best (L1) /\nGood (L2) thresholds editable\nby admin — saved to backend"),
    (9.6,  2.8, CYAN,   "Device Configuration",
     "Set OT type (OT / ICU) and\nroom name — used in CSV\nexport filenames"),
    (9.6,  4.3, GREEN,  "User Management",
     "Admin-only section: list users,\nadd with role, delete — persisted\nto users.json between restarts"),
    (9.6,  5.75, RED,   "Role-Based UI",
     "Save button hidden for viewers.\nUser Management section completely\nhidden for non-admin sessions"),
]

for x, y, col, title, body in callouts:
    rect(s, x, y, 3.5, 1.35, SURFACE)
    rect(s, x, y, 0.04, 1.35, col)
    txt(s, title, x+0.15, y+0.06, 3.2, 0.3, size=11, bold=True, color=col)
    txt(s, body,  x+0.15, y+0.38, 3.2, 0.9, size=9.5, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Dashboard: Alarms & History
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Dashboard — Alarms & History Panels",
             "Alarm log and 24-hour telemetry history pulled live from the SQLite database")

# Alarms
if os.path.exists(SHOTS["alarms"]):
    img(s, SHOTS["alarms"], 0.35, 1.1, 6.1, 3.8)
    border = s.shapes.add_shape(1, Inches(0.35), Inches(1.1), Inches(6.1), Inches(3.8))
    border.fill.background(); border.line.color.rgb = BORDER; border.line.width = Pt(1)

txt(s, "ALARM LOG", 0.35, 5.05, 4, 0.35, size=12, bold=True, color=AMBER)
for line in [
    "▸  Each entry: parameter name, timestamp, OT ID, value, level badge",
    "▸  Levels: NORMAL · WARNING · ALARM — colour coded",
    "▸  Pulls from GET /alarms?limit=100 on every navigation",
    "▸  Admin can acknowledge alarms via POST /alarms/{id}/acknowledge",
]:
    pass

alarm_pts = [
    "Parameter name + timestamp + OT ID",
    "WARNING / ALARM level badge (colour)",
    "Sorted latest-first from SQLite",
    "Auto-refreshes on nav click",
]
for i, pt in enumerate(alarm_pts):
    txt(s, "▸  " + pt, 0.45, 5.35 + i * 0.35, 5.9, 0.32, size=10.5, color=TEXT2)

# History
if os.path.exists(SHOTS["history"]):
    img(s, SHOTS["history"], 6.85, 1.1, 6.15, 3.8)
    border = s.shapes.add_shape(1, Inches(6.85), Inches(1.1), Inches(6.15), Inches(3.8))
    border.fill.background(); border.line.color.rgb = BORDER; border.line.width = Pt(1)

txt(s, "HISTORY TABLE", 6.85, 5.05, 4, 0.35, size=12, bold=True, color=CYAN)
hist_pts = [
    "Last 24 hours of telemetry rows",
    "Columns: Time / Temp / Humidity / PM2.5 / CO₂",
    "Pulls from GET /history?hours=24",
    "CSV export via DOWNLOAD button (auth required)",
]
for i, pt in enumerate(hist_pts):
    txt(s, "▸  " + pt, 6.95, 5.35 + i * 0.35, 6.0, 0.32, size=10.5, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Authentication & Security
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Authentication & Role-Based Access",
             "HMAC-SHA256 signed session cookies · Two roles: admin and viewer")

# Auth flow diagram
flow_steps = [
    (0.5,  2.8, CYAN,  "1  LOGIN",        "User clicks AUTH button\nEnters username + password"),
    (3.4,  2.8, AMBER, "2  VERIFY",        "POST /auth/login\nBackend checks users.json\n+ config.yaml"),
    (6.3,  2.8, GREEN, "3  SESSION COOKIE","HMAC-SHA256 token set\nas HttpOnly cookie\n(60 min expiry)"),
    (9.2,  2.8, CYAN,  "4  ROLE UI",       "Admin sees full controls\nViewer sees read-only\nUser Mgmt hidden"),
]

for x, y, col, title, body in flow_steps:
    rect(s, x, y, 2.6, 2.0, SURFACE)
    rect(s, x, y, 2.6, 0.05, col)
    txt(s, title, x+0.15, y+0.12, 2.3, 0.38, size=12, bold=True, color=col)
    hline(s, x+0.15, y+0.5, 2.3, BORDER)
    txt(s, body,  x+0.15, y+0.6, 2.3, 1.2, size=10.5, color=TEXT2)

# Arrow between steps
for ax in [3.15, 6.05, 8.95]:
    txt(s, "→", ax, 3.55, 0.4, 0.45, size=20, bold=True, color=TEXT2, align=PP_ALIGN.CENTER)

# Role table
rect(s, 0.5, 5.05, 12.3, 2.2, SURFACE)
txt(s, "Role Capabilities", 0.7, 5.12, 5, 0.38, size=12, bold=True, color=WHITE)
hline(s, 0.7, 5.5, 11.9, BORDER)

roles = [
    ("Capability",              "Admin 🔑",  "Viewer 👁"),
    ("Monitor panel (live data)", "✓",       "✓"),
    ("Download CSV",              "✓",       "✓"),
    ("View Alarms / History",     "✓",       "✓"),
    ("Edit thresholds",           "✓",       "✗  (hidden)"),
    ("Manage users",              "✓",       "✗  (hidden)"),
    ("Acknowledge alarms",        "✓",       "✗"),
]

col_x = [0.7, 7.5, 10.2]
colors_row = [TEXT2, GREEN, AMBER]
for j, (cap, adm, view) in enumerate(roles):
    y = 5.55 + j * 0.25
    bg = SURFACE2 if j % 2 == 1 else SURFACE
    rect(s, 0.5, y, 12.3, 0.25, bg)
    for ci, (val, col) in enumerate(zip([cap, adm, view], [WHITE, GREEN, AMBER])):
        is_header = j == 0
        txt(s, val, col_x[ci], y+0.03, 2.5, 0.22,
            size=9.5, bold=is_header, color=TEXT2 if is_header else col)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Alarm Engine
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Alarm Engine — Threshold Rules",
             "NABH / ASHRAE 170 compliant thresholds with hysteresis to prevent alarm fatigue")

# Threshold table
rect(s, 0.35, 1.15, 8.3, 5.9, SURFACE)
txt(s, "Parameter Thresholds (ASHRAE 170 / WHO 2021)", 0.55, 1.22, 7.5, 0.38,
    size=12, bold=True, color=WHITE)
hline(s, 0.55, 1.6, 8.0, BORDER)

headers = ["Parameter",  "Nominal",   "Warning Low", "Warning High", "Alarm Low",  "Alarm High"]
col_xs  = [0.55,          2.55,        3.75,           5.0,            6.2,           7.35]
col_w   = [1.9,           1.1,         1.15,           1.1,            1.1,           1.0]

for ci, h in enumerate(headers):
    txt(s, h, col_xs[ci], 1.65, col_w[ci], 0.3, size=9, bold=True, color=TEXT2)

thresholds = [
    ("Temperature (°C)",        "21",    "20",   "24",    "18",  "25",   CYAN),
    ("Humidity (%RH)",           "52",    "40",   "60",    "30",  "70",   CYAN),
    ("CO₂ (ppm)",               "—",     "—",    "1000",  "—",   "1200", AMBER),
    ("PM2.5 (µg/m³)",           "5",     "—",    "12",    "—",   "25",   RED),
    ("PM10 (µg/m³)",            "10",    "—",    "20",    "—",   "50",   RED),
    ("Pressure (hPa)",          "1013",  "990",  "1030",  "970", "1050", TEXT2),
]

for ri, (param, nom, wl, wh, al, ah, col) in enumerate(thresholds):
    y = 1.98 + ri * 0.78
    bg = SURFACE2 if ri % 2 == 0 else SURFACE
    rect(s, 0.55, y, 7.9, 0.72, bg)
    vals = [param, nom, wl, wh, al, ah]
    for ci, (v, cx, cw) in enumerate(zip(vals, col_xs, col_w)):
        vc = col if ci > 0 else WHITE
        txt(s, v, cx, y+0.2, cw, 0.38, size=10, color=vc)

# Hysteresis explanation
rect(s, 8.85, 1.15, 4.1, 5.9, SURFACE)
txt(s, "Hysteresis Design", 9.05, 1.22, 3.7, 0.38, size=12, bold=True, color=WHITE)
hline(s, 9.05, 1.6, 3.7, BORDER)

hyst_points = [
    (AMBER, "Trigger Delay",
     "Parameter must exceed\nthreshold continuously\nfor 5 seconds before\nan alarm fires"),
    (GREEN, "Clear Delay",
     "Must stay back in safe\nrange for 10 seconds\nbefore alarm clears\n(prevents chattering)"),
    (CYAN,  "Combination Rule",
     "PM2.5 spike + high\noccupancy together\ntrigger elevated-risk\nnotification"),
    (TEXT2, "Alarm Levels",
     "NORMAL → WARNING\n→ ALARM → FAULT\nEach stored in\nSQLite alarm log"),
]

for i, (col, title, body) in enumerate(hyst_points):
    y = 1.7 + i * 1.35
    rect(s, 9.05, y, 3.7, 1.25, SURFACE2)
    rect(s, 9.05, y, 0.04, 1.25, col)
    txt(s, title, 9.2, y+0.07, 3.4, 0.3, size=10, bold=True, color=col)
    txt(s, body,  9.2, y+0.38, 3.4, 0.8, size=9.5, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Data Flow
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Data Flow — End to End",
             "From physical sensor to dashboard KPI card in under 2 seconds")

steps = [
    (0.35, CYAN,   "SENSORS",        "SCD30 + BME280 + PMS5003\nRead every 2 seconds\nvia asyncio thread pool"),
    (2.65, AMBER,  "HARDWARE SOURCE", "HardwareSource assembles\nTelemetryPacket\nTimestamped in IST"),
    (4.95, GREEN,  "ALARM ENGINE",   "Evaluate thresholds\nApply hysteresis rules\nFire callbacks on breach"),
    (7.25, CYAN,   "STORAGE",        "Insert telemetry row\nInsert alarm event\nSQLite WAL mode"),
    (9.55, RED,    "REST API",       "GET /history\nGET /alarms\nWebSocket broadcast"),
    (11.85,WHITE,  "DASHBOARD",      "KPI cards update\nAlarm banners show\nHistory table fills"),
]

for x, col, title, body in steps:
    rect(s, x, 1.3, 2.1, 2.5, SURFACE)
    rect(s, x, 1.3, 2.1, 0.05, col)
    txt(s, title, x+0.12, 1.38, 1.9, 0.38, size=9, bold=True, color=col)
    hline(s, x+0.12, 1.76, 1.86, BORDER)
    txt(s, body, x+0.12, 1.85, 1.86, 1.8, size=9.5, color=TEXT2)

# Arrows
for ax in [2.5, 4.8, 7.1, 9.4, 11.7]:
    txt(s, "▶", ax, 2.25, 0.3, 0.4, size=14, color=BORDER, align=PP_ALIGN.CENTER)

# Timing bar
hline(s, 0.35, 4.1, 12.6, BORDER)
txt(s, "Total end-to-end latency: < 2 seconds (poll interval = 2 s)",
    0.35, 4.15, 12.6, 0.38, size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Storage detail
rect(s, 0.35, 4.65, 12.6, 2.55, SURFACE)
txt(s, "Storage & Retention", 0.55, 4.72, 5, 0.38, size=12, bold=True, color=WHITE)
hline(s, 0.55, 5.1, 12.2, BORDER)

storage_info = [
    ("Database",       "SQLite with WAL mode — concurrent reads without locking the write path"),
    ("Telemetry",      "1 row per 2-second poll · fields: ot_id, timestamp, temp, humidity, PM2.5, CO₂, pressure, status"),
    ("Alarms",         "Every threshold breach logged: parameter, level, value, message, acknowledged, duration"),
    ("Retention",      "90 days default (configurable) · auto-pruned every 6 hours by background task"),
    ("Export",         "GET /export/csv → downloadable audit report for NABH / infection control review"),
]

for i, (k, v) in enumerate(storage_info):
    y = 5.15 + i * 0.38
    txt(s, k + ":", 0.55, y, 1.5,  0.32, size=10, bold=True, color=CYAN)
    txt(s, v,       2.1,  y, 10.6, 0.32, size=10, color=TEXT2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — RPi Deployment
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Raspberry Pi Deployment",
             "From fresh SD card to live dashboard in under 30 minutes")

steps_deploy = [
    ("1", CYAN,  "Flash SD Card",
     "Raspberry Pi Imager → RPi OS Lite 64-bit\nPre-configure SSH, hostname: ot-monitor"),
    ("2", AMBER, "Enable Interfaces",
     "/boot/firmware/config.txt:\n  dtparam=i2c_arm_baudrate=10000\n  dtoverlay=disable-bt\n  enable_uart=1"),
    ("3", GREEN, "Install Libraries",
     "pip install adafruit-circuitpython-scd30\nadafruit-circuitpython-bme280 pyserial"),
    ("4", CYAN,  "Test Sensors",
     "python sensors/test_sensors.py\nAll three sensors must show PASS\nbefore going to hardware mode"),
    ("5", RED,   "Configure & Run",
     "config.yaml: data_source.type: hardware\nserver.host: 0.0.0.0\nChange session_secret + passwords"),
    ("6", WHITE, "Autostart",
     "systemd service: ot-monitor.service\nsudo systemctl enable ot-monitor\nDashboard at :8001 on every boot"),
]

for i, (num, col, title, body) in enumerate(steps_deploy):
    row = i // 3
    col_idx = i % 3
    x = 0.35 + col_idx * 4.3
    y = 1.15 + row * 2.9
    rect(s, x, y, 4.1, 2.65, SURFACE)
    rect(s, x, y, 4.1, 0.05, col)
    # Step number
    rect(s, x+0.15, y+0.12, 0.4, 0.4, col)
    txt(s, num, x+0.15, y+0.12, 0.4, 0.4, size=14, bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, title, x+0.65, y+0.14, 3.3, 0.38, size=13, bold=True, color=col)
    hline(s, x+0.15, y+0.55, 3.8, BORDER)
    txt(s, body, x+0.15, y+0.65, 3.8, 1.8, size=10.5, color=TEXT2)

txt(s, "Access dashboard:  http://ot-monitor.local:8001/   or   http://<RPi-IP>:8001/",
    0.35, 6.95, 12.6, 0.38, size=12, bold=True, color=CYAN, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Clinical Standards
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
slide_header(s, "Clinical & Regulatory Compliance",
             "Thresholds derived from NABH, ASHRAE 170-2021, WHO IEQ 2021, ISO 14644-1")

standards = [
    (CYAN,  "NABH",
     "National Accreditation Board for Hospitals & Healthcare Providers (India)\n"
     "OT environment: 20–24°C, 40–60%RH, positive pressure ≥ 8 Pa,\n"
     "≥ 20 Air Changes per Hour (ACH), HEPA filtration."),
    (AMBER, "ASHRAE 170-2021",
     "Ventilation of Health Care Facilities — North American standard adopted globally.\n"
     "OT: 20–24°C, 30–60%RH, ≥20 ACH total / ≥4 ACH outside air,\n"
     "PM2.5 ≤ 25 µg/m³, PM10 ≤ 50 µg/m³."),
    (GREEN, "WHO IEQ 2021",
     "Indoor Air Quality Guidelines (updated 2021).\n"
     "PM2.5 24-hr average ≤ 15 µg/m³, PM10 ≤ 45 µg/m³,\n"
     "CO₂ < 1000 ppm in occupied clinical spaces."),
    (RED,   "ISO 14644-1",
     "Cleanrooms and Associated Controlled Environments.\n"
     "OT typically ISO Class 7 (Class 10,000):\n"
     "≤ 352,000 particles ≥ 0.5µm per m³ — monitored via PMS5003."),
]

for i, (col, std, body) in enumerate(standards):
    row = i // 2
    ci  = i % 2
    x = 0.35 + ci * 6.5
    y = 1.15 + row * 2.85
    rect(s, x, y, 6.2, 2.6, SURFACE)
    rect(s, x, y, 6.2, 0.05, col)
    txt(s, std,  x+0.2, y+0.12, 5.8, 0.42, size=16, bold=True, color=col)
    hline(s, x+0.2, y+0.55, 5.8, BORDER)
    txt(s, body, x+0.2, y+0.65, 5.8, 1.75, size=10.5, color=TEXT2)

txt(s, "All thresholds stored in config/config.yaml — editable via Settings panel (admin only) without restarting the backend",
    0.35, 6.95, 12.6, 0.38, size=11, color=TEXT2, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Summary / Thank You
# ═══════════════════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 0.08, CYAN)
rect(s, 0, 0.08, 0.05, 7.42, CYAN)

txt(s, "SYSTEM SUMMARY", 0.6, 0.8, 10, 0.55, size=32, bold=True, color=WHITE)
hline(s, 0.6, 1.45, 5.5, CYAN, 2)

features = [
    (CYAN,  "Complete Hardware Stack",
     "Raspberry Pi 4 + SCD30 + BME280 + PMS5003 — all sensors integrated and tested"),
    (GREEN, "Real-Time Monitoring",
     "6 environmental parameters polled every 2 seconds with live dashboard updates"),
    (AMBER, "Clinical Alarm Engine",
     "NABH / ASHRAE 170 thresholds with hysteresis — zero false alarm fatigue"),
    (RED,   "Secure Multi-User Access",
     "HMAC-SHA256 sessions, role-based UI — admin vs viewer from any browser"),
    (CYAN,  "90-Day Data Retention",
     "SQLite time-series + alarm log, CSV export for audit and compliance reporting"),
    (WHITE, "Zero-Dependency Dashboard",
     "Single HTML file served from backend — no build step, works on any device"),
]

for i, (col, title, desc) in enumerate(features):
    y = 1.65 + i * 0.82
    rect(s, 0.6, y, 0.06, 0.55, col)
    txt(s, title, 0.85, y,      6.5, 0.32, size=13, bold=True, color=col)
    txt(s, desc,  0.85, y+0.3, 6.5, 0.35, size=10.5, color=TEXT2)

# Right side — logout screenshot
if os.path.exists(SHOTS["logout"]):
    img(s, SHOTS["logout"], 7.6, 1.35, 5.5, 3.5)
    border = s.shapes.add_shape(1, Inches(7.6), Inches(1.35), Inches(5.5), Inches(3.5))
    border.fill.background(); border.line.color.rgb = CYAN; border.line.width = Pt(1.5)

rect(s, 7.6, 5.0, 5.5, 1.0, SURFACE)
txt(s, "http://ot-monitor.local:8001/",
    7.6, 5.05, 5.5, 0.42, size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
txt(s, "admin / OTAdmin2024  ·  nurse / OTNurse2024",
    7.6, 5.45, 5.5, 0.38, size=10, color=TEXT2, align=PP_ALIGN.CENTER)

txt(s, "OT Infection Monitoring System  ·  Raspberry Pi 4 + SCD30 + BME280 + PMS5003  ·  June 2026",
    0.4, 7.05, 12.5, 0.35, size=9, color=BORDER, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = os.path.join(os.path.dirname(__file__), "OT_Monitor_Presentation.pptx")
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
